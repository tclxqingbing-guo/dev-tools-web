import os
from pathlib import Path


def read_document(path_value: str) -> str:
    """解析会话附件为纯文本，不执行附件中的宏或脚本。"""
    path = Path(path_value).resolve()
    root = Path(os.environ.get('AGENT_ATTACHMENT_ROOT', '/app/data/agent/attachments')).resolve()
    if root not in path.parents or not path.is_file():
        raise ValueError('附件路径无效')
    suffix = path.suffix.lower()
    if suffix in {'.txt', '.md', '.markdown', '.csv', '.json', '.xml'}:
        return path.read_text('utf-8', errors='replace')[:200000]
    if suffix in {'.html', '.htm'}:
        from bs4 import BeautifulSoup
        return BeautifulSoup(path.read_text('utf-8', errors='replace'), 'html.parser').get_text('\n')[:200000]
    if suffix == '.pdf':
        from pypdf import PdfReader
        return '\n\n'.join((page.extract_text() or '') for page in PdfReader(str(path)).pages)[:200000]
    if suffix == '.docx':
        from docx import Document
        document = Document(str(path))
        return '\n'.join(p.text for p in document.paragraphs)[:200000]
    if suffix == '.pptx':
        from pptx import Presentation
        presentation = Presentation(str(path))
        return '\n'.join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, 'text'))[:200000]
    if suffix in {'.png', '.jpg', '.jpeg', '.webp', '.tif', '.tiff', '.bmp'}:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(path), lang='chi_sim+eng')[:200000]
    raise ValueError(f'暂不支持该附件格式: {suffix or "unknown"}')
